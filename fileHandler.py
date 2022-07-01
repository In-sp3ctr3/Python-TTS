import os
import glob
import PyPDF2
import docx2txt
import speechHandler as speechHandler
from PyPDF2 import PdfFileMerger
from pptx import Presentation
from nltk import Tree, tree
from nltk.draw.util import CanvasFrame
from nltk.draw import TreeWidget
from languageProcessing.lexicalAnalysisEngine.lexical_analyzer import Analyzer


def file_read(name_of_file, treeFileName):
    path = "uploads/" + name_of_file
    extension = name_of_file.split('.')[-1]
    name_of_file_alone = ( name_of_file.rsplit( ".", 1 )[ 0 ] )
    success = True
    text = ''
    
    try:
        if extension == 'txt':
            if is_file_empty(path):
                success = False
                raise Exception("File string is empty!")
            else:
                with open(path,"r") as f:
                    text = f.read()
        elif extension == 'pdf':
            pdfFileObj = open(path, "rb")
            pdfReader = PyPDF2.PdfFileReader(pdfFileObj)
            totalPages = pdfReader.numPages
            for i in range(totalPages):
                page = pdfReader.getPage(i)
                text += page.extractText()
        elif extension == 'docx':
            text = docx2txt.process(path)
        elif extension == 'pptx':
            pptOutput = ''
            ppt = Presentation(path)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text
    except:
        print("An error occurred!")
    return text

def is_file_empty(path):
    with open(path, 'r') as fileReader:
        first_char = fileReader.read(1)
        if not first_char:
            return True
    return False


def generate_tree_pdf(extractions, treeFileName):
    merger = PdfFileMerger()

    for x,extraction in enumerate(extractions):
        print(str(extraction))
        cf = CanvasFrame()
        t = Tree.fromstring(str(extraction))
        tc = TreeWidget(cf.canvas(),t)
        tc['node_font'] = 'arial 14 bold'
        tc['leaf_font'] = 'arial 14'
        tc['node_color'] = '#005990'
        tc['leaf_color'] = '#3F8F57'
        tc['line_color'] = '#175252'
        cf.add_widget(tc,10,10) 
        ps_file_name = 'tree'+"-"+str(x)+".ps"
        pdf_file_name = 'tree'+"-"+str(x)+".pdf"
        cf.print_to_file('parseDocs/' + ps_file_name)
        cf.destroy()
        os.system('magick convert parseDocs/' + ps_file_name + ' parseDocs/' + pdf_file_name)

        merger.append('parseDocs/' + pdf_file_name)
    
    merger.write('parseDocs/' + treeFileName)
    merger.close()
